#!/usr/bin/env python
# -*- coding: utf-8 -*-
# @Date  : 2020/12/22 10:34 上午
# @File  : use_api.py
# @Author: johnson
# @Contact : github: johnson7788
# @Desc  :  翻译的任务
"""
输入text是中文，需要的是翻译成英文
"""

import requests
import json
import pprint
import os
import hashlib
import time
import re
import zipfile
from pptx import Presentation

headers = {'content-type': 'application/json;charset=UTF-8'}
host = "http://localhost:8080/api/"
localhost = "http://localhost:8080/api/"
pp = pprint.PrettyPrinter(indent=4)

def setup_config(hostname=None):
    """
    ner分类的config
    :return:
    """
    data = {"label_config":
                """
<View>
  <View style="flex: 30%; color:red">
    <Text name="text" value="$text"/>
    <TextArea name="result"></TextArea>
    </View>
    <View style="flex: 30%">
      <Labels name="label" toName="text">
        <Label value="跳过翻译" background="blue"></Label>
      </Labels>
  </View>
</View>
                """}
    if hostname != None:
        host = hostname
    r = requests.post(host + "project/config", data=json.dumps(data), headers=headers)
    print(r.status_code)
    print(r.text)


def setup_config_host(hostnames):
    """
    删除所有task, 数据, 同时会删除已标注的数据
    :param hostnames: 要设置的hostname
    :return:
    """
    for hostname in hostnames:
        setup_config(hostname=hostname)


def get_project():
    """
    获取项目的信息
    :return:
    """
    r = requests.get(host + "project/", headers=headers)
    print(r.json())
    pp.pprint(r.json())


def get_tasks(taskid=None, page_size=5000, hostname=None):
    """
    获取task, 获取数据，如果标注完成，返回标注的状态
    :param: taskid 获取第几条数据，如果为None，获取所有数据
    :return:
    """
    if hostname != None:
        host = hostname
    if taskid:
        taskid = str(taskid)
        r = requests.get(host + "tasks/" + taskid, data=json.dumps({'filters':None}), headers=headers)
    else:
        payload = {'fields': 'all', 'page': 1, 'page_size': page_size, 'order': 'id'}
        r = requests.get(host + "tasks", params=payload, data=json.dumps({'filters':None}), headers=headers)
    # print(r.json())
    # pp.pprint(r.json())
    results = r.json()
    return results


def get_tasks_host(hostnames):
    """
    删除所有task, 数据, 同时会删除已标注的数据
    :return:
    """
    for hostname in hostnames:
        get_tasks(hostname=hostname)


def delete_tasks(hostname, taskid=None):
    """
    如果taskid=None删除所有task, 数据, 同时会删除已标注的数据
    :return:
    """
    if taskid:
        taskid = str(taskid)
        r = requests.delete(hostname + "tasks/" + taskid, headers=headers)
    else:
        r = requests.delete(hostname + "tasks", headers=headers)
    print("完成，返回code:")
    print(r.status_code)
    print(r.text)


def delete_tasks_host(hostnames):
    """
    删除所有task, 数据, 同时会删除已标注的数据
    :return:
    """
    for hostname in hostnames:
        delete_tasks(hostname)


def get_completions(taskid=None, hostname=None, proxy=False):
    """
    获取完成的task，默认获取所有完成的样本，可以获取部分样本,指定taskid
    :param taskid:
    :return:  返回完成的task的id， task就是样本， {"ids":[0,1]}
    """
    if hostname != None:
        host = hostname
    if taskid:
        taskid = str(taskid)
        if proxy:
            r = requests.get(host + "tasks/" + taskid + "/completions", headers=headers,
                             proxies=dict(http='socks5://127.0.0.1:9080', https='socks5://127.0.0.1:9080'))
        else:
            r = requests.get(host + "tasks/" + taskid + "/completions", headers=headers)
    else:
        if proxy:
            r = requests.get(host + "completions", headers=headers,
                             proxies=dict(http='socks5://127.0.0.1:9080', https='socks5://127.0.0.1:9080'))
        else:
            r = requests.get(host + "completions", headers=headers)
    # print(r.status_code)
    # print(r.text)
    complete_ids = r.json()["ids"]
    print(f"{host}  标注完成{len(complete_ids)}条")
    return len(complete_ids)


def get_completions_host(hostnames):
    """
    删除所有task, 数据, 同时会删除已标注的数据
    :return:
    """
    print(time.strftime("%Y/%m/%d %H:%M:%S", time.localtime()))
    total = 0
    for hostname in hostnames:
        complete_num = get_completions(hostname=hostname)
        total += complete_num
    print(f"总共已完成{total}条数据标注")


def delete_completions():
    """
    删除所有已完成的标注样本
    :return:
    """
    r = requests.delete(host + "completions", headers=headers)
    print(r.status_code)
    print(r.text)


def import_data(hostname):
    """
    导入字典里面包含多个key和value的格式
    例如
    data = [{"text": "很好，实惠方便，会推荐朋友", "channel":"jd", "keyword":""},{"text": "一直买的他家这款洗发膏，用的挺好的，洗的干净也没有头皮屑"}]
    channel	content	关键词	判断是否准确	理由
tmall	有薄荷，抹上凉凉的，没有什么不适	薄荷	准确	说的是成分
tmall	整体评价：很大一瓶 水的质地比较浓稠 有股那种洗护的味道 刺鼻 不好闻 上脸比较滋润 鼻翼附近会刺痛 一分钱一分货吧 感受不好 不会回购 我一直在用红石榴和微精华水这两款水肤感都很好 推荐购买	红石榴	不准确	红石榴指的是产品
tmall	发货特别的快，我第1天买的，第2天就到了，迫不及待的打开试用了一下，吸收效果不错的，提亮肤色效果特别好，我的肤质比较偏干，正好适合秋冬季用，只是味道有一点点不太喜欢，自己加了两滴玫瑰精油一起，还挺好的。本来以为送了一盒膜，打开发现是个祝福的相框吧，还挺尴尬的。哈哈???	玫瑰	不准确	说的是精油，并非成分
tmall	不错，湿敷比擦拭更舒服，冰冰凉凉的有点提神，好闻的薄荷感	薄荷	不准确	说的是使用感受，并非成分
tmall	买来收敛毛孔的，看着是天猫超市，价格也便宜，标签包装都挺完整的，我都一层一层撕下来了，闻着是橙子和牙膏薄荷的味道。这个的产品有好多假冒的，希望是真的吧！	薄荷	不准确	说的是味道和使用感受，并非成分
tmall	湿敷的，有薄荷跟橘子的味道，才用，期待以后的效果吧	薄荷	不准确	说的是味道，并非成分
tmall	这个水可以用来湿敷 虽然用来湿敷有些贵 ！里面含有薄荷，涂上非常舒服，洗面奶面霜也炒鸡适合我 我油皮痘痘肌	薄荷	准确	说的是成分
redbook	今天分享一下我在夏天非常喜欢用的三款水：1??悦木之源菌菇水真的太太太喜欢啦，非常的清爽，不油腻，上脸吸收速度很快，湿敷对痘痘消肿有不错的效果。常年囤货。2??skii神仙水控油效果非常明显，对油皮和混油皮友好，干皮渗入。除了贵和味道之外，没有别的缺点，有预算的宝宝买就对了。3??城野医生收敛水可以调理你的皮脂分泌，但是没有收缩毛孔和去黑头的功效，使用感很好，清清凉凉的，但其实不添加酒精，而是添加了薄荷醇，肤感和气味都和酒精有点类似，但是很温和，不会造成屏障受损。@生活薯  @薯队长	薄荷	准确	说的是成分
tmall	用在白泥之后还是很不错的，觉得有二次湿敷出白头，收缩毛孔大概是有的吧	白泥	不准确	说的是产品，并非成分
weibo	#空瓶记# #空瓶记城野医生的毛孔收敛水 每周毛孔清洁后必备的水儿 配合奥尔滨的化妆棉进行二次清洁 平价好用 推了一万年了 我去日本就背了七八瓶回来颐莲的玻尿酸护发原液 因为我头发经常烫染所以发梢有点受损 每次洗完头发按三泵涂抹 不粘腻及时补水然后再涂抹护发精油锁住水份 颐莲这个牌子是福瑞达旗下的 福瑞达和华西生物都是国内做玻尿酸的大牛 雅诗兰黛的玻尿酸供应商 所以不要瞧不起他们 有这么便宜好用性价比高的国货不支持 还等啥呢 收起全文d	玻尿酸	不准确	说的是产品，并非成分

    :return:
    """
    data = [
        {'text': '国货护肤专题研究'},
        {'text': '针对雅诗的建议'},
        {'text': '国货护肤的现状：'},
        {'text': '2018年起，中国开始进入国潮热，很多品牌开始借势营销，自2020年起，品牌更注重内功，用功效和品质圈定年轻消费者。'},
        {'text': '国货护肤主要受众人（Z世代）的特征：'},
        ]
    r = requests.post(hostname + "project/import", data=json.dumps(data), headers=headers)
    pp.pprint(r.json())


def health():
    """
    测试健康
    :return:
    """
    r = requests.get(host + "health", headers=headers)
    print(r.status_code)
    print(r.text)
    r = requests.get("http://localhost:8080/version", headers=headers)
    print(r.status_code)
    print(r.text)


def list_models():
    """
    模型后端api
    :return:
    """
    r = requests.get(host + "models", headers=headers)
    print(r.status_code)
    print(r.text)


def train_model():
    """
    训练模型, 训练所有的后端模型, 只有在有新的标注后，才调用后端自定义的fit接口，没有新的标注，就不调用后端的fit训练模型
    例如：调用absa_classifier.py的ABSATextClassifier类的init，然后调用fit
    label_studio 调用 label_studio_ml的  'http://localhost:9090/train'接口
    调用的post的参数
    request = {dict: 4} {'completions': [{'completions': [{'created_at': 1608617962, 'id': 1, 'lead_time': 5.162, 'result': [{'from_name': 'sentiment', 'id': '8BAF02Fcq5', 'to_name': 'text', 'type': 'choices', 'value': {'choices': ['积极']}}]}], 'data': {'text': '很好，实惠方便，会推荐朋友'}, '
     'completions' = {list: 4} [{'completions': [{'created_at': 1608617962, 'id': 1, 'lead_time': 5.162, 'result': [{'from_name': 'sentiment', 'id': '8BAF02Fcq5', 'to_name': 'text', 'type': 'choices', 'value': {'choices': ['积极']}}]}], 'data': {'text': '很好，实惠方便，会推荐朋友'}, 'id': 0}, {'completions': [{'created_at': 1608617965, 'id': 1001, 'lead_time': 3.101, 'result': [{'from_name': 'sentiment', 'id': 'UvcGY6ZCM7', 'to_name': 'text', 'type': 'choices', 'value': {'choices': ['中性']}}]}], 'data': {'text': '一直买的他家这款洗发膏，用的挺好的，洗的干净也没有头皮屑'}, 'id': 1}, {'completions': [{'created_at': 1608617968, 'id': 2001, 'lead_time': 2.085, 'result': [{'from_name': 'sentiment', 'id': 'h1qNhxHSyQ', 'to_name': 'text', 'type': 'choices', 'value': {'choices': ['消极']}}]}], 'data': {'text': '不太顺滑'}, 'id': 2}, {'completions': [{'created_at': 1608617970, 'id': 3001, 'lead_time': 2.583, 'result': [{'from_name': 'sentiment', 'id': 'V0Q1xEHQQX', 'to_name': 'text', 'type': 'choices', 'value': {'choices': ['积极']}}]}], 'data': {'text': '特别香，持久'}, 'id': 3}]
     'project' = {str} 'text_classification_project1a43'
     'label_config' = {str} '<View>  <Text name="text" value="$text"/>  <Choices name="sentiment" toName="text" choice="single">    <Choice value="积极"/>    <Choice value="消极"/>    <Choice value="中性"/>  </Choices></View>'
     'params' = {dict: 3} {'login': '', 'password': '', 'project_full_path': '/Users/admin/git/label-studio/text_classification_project'}
     __len__ = {int} 4
    :return:
    """
    r = requests.post(host + "models/train", headers=headers)
    print(r.status_code)
    print(r.text)


def predict_model():
    """
    调用模型的预测接口, 首先调用自定义的分了器，例如absa_classifier.py的init，然后对所有tasks进行调用absa_classifier.py的predict进行预测
    /api/models/predictions?mode={data|all_tasks}
    :return: eg: {"details":"predictions are ready"}
    """
    r = requests.post(host + "models/predictions?mode=all_tasks", headers=headers)
    print(r.status_code)
    print(r.text)


def cal_md5(content):
    """
    计算content字符串的md5
    :param content:
    :return:
    """
    # 使用encode
    result = hashlib.md5(content.encode())
    # 打印hash
    md5 = result.hexdigest()
    return md5


def get_imported_data_md5(imported_data):
    """
    对已经导入的数据，计算所有md5，如果data['md5']存在，直接过去，否则用keyword+text计算md5
    :param: imported_data是每个tasks
    :return: 按列表顺序返回md5的列表[]
    """
    md5_list = []
    for res in imported_data:
        data = res['data']
        md5_value = data.get('md5')
        if not md5_value:
            # 说明不存在md5这个字段，开始计算
            content = data['keyword'] + data['text']
            md5_value = cal_md5(content=content)
        md5_list.append(md5_value)
    return md5_list


def import_com_data_host(channel=['jd', 'tmall'], number=10, hostname=None):
    """
    按比例导入不同的host, 导入情感分析数据, 从hive数据库中导入, 导入到label-studio前，需要检查下这条数据是否已经导入过
    12月份，功效4000条，其它维度各1500条
    :param number:
    :param hostname:平均导入每个host中,列表或None
    :return:
    """
    leibie = ['成分', '功效', '香味', '包装', '肤感']
    from read_hive import get_absa_corpus
    # 要导入的数据
    valid_data = []
    # 已经导入的数据, 注意更改获取的样本数目，默认是5000条
    if hostname != None:
        host = hostname
    imported_data = []
    for h in host:
        host_imported_data = get_tasks(page_size=5000, hostname=h)
        imported_data.extend(host_imported_data)
    imported_data_md5 = get_imported_data_md5(imported_data)
    # 开始从hive数据库拉数据
    data = get_absa_corpus(channel=channel, requiretags=None, number=number)
    # 获取到的data数据进行排查，如果已经导入过了，就过滤掉
    for one_data in data:
        content = one_data['keyword'] + one_data['text']
        data_md5 = cal_md5(content)
        if data_md5 in imported_data_md5:
            # 数据已经导入到label-studio过了，不需要重新导入
            continue
        else:
            # 没有导入过label-studio，那么加入到valid_data，进行导入
            # 设置md5字段，方便以后获取
            one_data['md5'] = data_md5
            valid_data.append(one_data)
    print(f"可导入的有效数据是{len(valid_data)}, 有重复数据{len(data) - len(valid_data)} 是无需导入的")
    if not valid_data:
        # 如果都是已经导入过的数据，直接放弃导入
        return
    every_host_number = int(len(valid_data) / len(host))
    print(f"每个主机导入数据{every_host_number}")
    vdatas = [valid_data[i:i + every_host_number] for i in range(0, len(valid_data), every_host_number)]
    for h, vdata in zip(host, vdatas):
        r = requests.post(h + "project/import", data=json.dumps(vdata), headers=headers)
        pp.pprint(r.json())
        print(f"共导入主机host{h}中数据{len(vdata)}条")


def import_com_data_host_first(channel=['jd', 'tmall'], keyword_threhold=0, leibie_num=[1000, -1, -1, -1, -1], require_tags=["component","effect","fragrance","pack","skin"], number=10, unique_type=1, hostname=None, ptime_keyword=">:2020-12-01", not_cache=True, table="da_wide_table_new"):
    """
    按比例导入不同的host, 导入成分分析数据, 从hive数据库中导入, 导入到label-studio前，需要检查下这条数据是否已经导入过
    12月份，功效4000条，其它维度各1500条
    :param channel: channel:包括 "jd","weibo","redbook","tiktok","tmall"
    :param keyword_threhold: 0表示不进行keyword的过滤，如果给定数字，例如10，表示这个keyword不会在过滤次数中超过10次
    :param number: 从数据库检索的条目，不是最终数目，最终数目根据leibie_num确定
    :param require_tags: 需要哪些维度的语料
    :param hostname:平均导入每个host中,列表或None
    :return:
    """
    from math_tags import get_emotional_words
    confirm_file = "com_100.txt"
    all_tags = []
    #收集所有tags
    with open(confirm_file, 'r') as f:
        for line in f:
            line = line.strip()
            tmp = {}
            tmp['0'] = line.split('|')
            all_tags.append(str(tmp))
    leibie = ['成分', '功效', '香味', '包装', '肤感']
    from read_hive import get_absa_corpus
    # 要导入的数据
    valid_data = []
    # 已经导入的数据, 注意更改获取的样本数目，默认是5000条
    if hostname != None:
        host = hostname
    imported_data = []
    for h in host:
        print(f"获取{h}的任务")
        host_imported_data = get_tasks(page_size=8000, hostname=h)
        imported_data.extend(host_imported_data['tasks'])
    imported_data_md5 = get_imported_data_md5(imported_data)
    # 开始从hive数据库拉数据, 如果unique_type设置为2，那么数据可能过少
    data = get_absa_corpus(channel=channel, requiretags=require_tags, number=number, unique_type=unique_type, ptime_keyword=ptime_keyword, not_cache=not_cache, table=table)
    # 获取到的data数据进行排查，如果已经导入过了，就过滤掉
    initial_count = [0, 0, 0, 0, 0]
    #keywords的unique记录
    keywords_dict = {}
    for one_data in data:
        get_index = leibie.index(one_data['wordtype'])
        if initial_count[get_index] < leibie_num[get_index]:
            initial_count[get_index] += 1
        else:
            continue
        content = one_data['keyword'] + one_data['text']
        data_md5 = cal_md5(content)
        #用于匹配当前这个数据的关键字是否和我们提供的100%确定是成分的关键字匹配
        match_keywords = []
        #已经获取了多少个这个关键字的示例, 如果不存在，那么为0
        keyword_num = keywords_dict.get(one_data['keyword'], 0)
        for current_tags in all_tags:
            keywords, keyword_dictid = get_emotional_words(tag_words=current_tags, content=content)
            match_keywords.extend(keywords)
        if data_md5 in imported_data_md5:
            # 数据已经导入到label-studio过了，不需要重新导入
            continue
        if one_data['keyword'] in match_keywords:
            print(f'这个句子的关键字: {one_data["keyword"]} 在白名单中出现，跳过，句子是: {one_data["text"]}')
            continue
        elif keyword_threhold !=0 and keyword_num > keyword_threhold:
            #这个keyword出现的次数已经超过所需要阈值，可以过滤掉, 加到数据库中，统计下
            keywords_dict[one_data['keyword']] = keyword_num + 1
            continue
        else:
            # 没有导入过label-studio，那么加入到valid_data，进行导入
            # 设置md5字段，方便以后获取
            # 把keywords中的这个关键字的数量+1
            keywords_dict[one_data['keyword']] = keyword_num + 1
            one_data['md5'] = data_md5
            valid_data.append(one_data)
    print(f'关键字出现的总的次数：{keywords_dict}')
    print(f"可导入的有效数据是{len(valid_data)}, 有重复数据或不需要数据{len(data) - len(valid_data)} 是无需导入的")
    if not valid_data:
        # 如果都是已经导入过的数据，直接放弃导入
        return
    if len(valid_data) < number:
        print(f"收集的数据量过少，很可能是因为get_absa_corpus的unique type设置问题")
    every_host_number = int(len(valid_data) / len(host))
    print(f"每个主机导入数据{every_host_number}")
    vdatas = [valid_data[i:i + every_host_number] for i in range(0, len(valid_data), every_host_number)]
    for h, vdata in zip(host, vdatas):
        r = requests.post(h + "project/import", data=json.dumps(vdata), headers=headers)
        pp.pprint(r.json())
        print(f"共导入主机host: {h}中数据{len(vdata)}条")


def check_data(hostname):
    """
    查看下已导入的数据
    :return:
    """
    datas = get_tasks(hostname)
    not_repeat_id = []
    not_repeat_data = []
    for data in datas:
        content = data['data']['keyword'] + data['data']['text']
        if content not in not_repeat_data:
            not_repeat_data.append(content)
            not_repeat_id.append(data['id'])
        else:
            repeat_idx = not_repeat_data.index(content)
            repeat_id = not_repeat_id[repeat_idx]
            print(f"发现重复数据:{data['id']}和{repeat_id}")
    print(f"共有重复数据{len(datas) - len(not_repeat_data)}条")


def export_data(hostname=None, dirpath="/opt/lavector/", jsonfile=None, proxy=False):
    """
    导出数据
    :param hostname:
    :return:
    """
    if hostname != None:
        host = hostname
    # 获取下标注完成了多少了数据
    get_completions(hostname=host, proxy=proxy)
    p = '(?:http.*://)?(?P<host>[^:/ ]+).?(?P<port>[0-9]*).*'
    m = re.search(p, host)
    hostip = m.group('host')
    port = m.group('port')
    url = host + "project/export?format=JSON"
    local_zipfile = hostip + "_" + port + "_" + time.strftime("%Y%m%d%H%M%S", time.localtime()) + ".zip"
    local_zipfile = os.path.join("/tmp", local_zipfile)
    if jsonfile is None:
        local_jsonfile = hostip + "_" + port + ".json"
        local_jsonfile = os.path.join(dirpath, local_jsonfile)
    else:
        local_jsonfile = os.path.join(dirpath, jsonfile)
    if os.path.exists(local_jsonfile):
        print(f"注意，json文件已经存在即将覆盖json文件 {local_jsonfile}")
    if proxy:
        with requests.get(url, stream=True,
                          proxies=dict(http='socks5://127.0.0.1:9080', https='socks5://127.0.0.1:9080')) as r:
            r.raise_for_status()
            with open(local_zipfile, 'wb') as f:
                for chunk in r.iter_content(chunk_size=8192):
                    f.write(chunk)
    else:
        with requests.get(url, stream=True) as r:
            r.raise_for_status()
            with open(local_zipfile, 'wb') as f:
                for chunk in r.iter_content(chunk_size=8192):
                    f.write(chunk)
    # 创建一个压缩包对象
    parent_archive = zipfile.ZipFile(local_zipfile)
    # 解压
    parent_archive.extractall(dirpath)
    # 获取文件夹的压缩包列表
    namelist = parent_archive.namelist()
    files = [dirpath + name for name in namelist]
    parent_archive.close()
    assert (len(files) == 1), "压缩包里面不止一个文件，请检查"
    extract_file = files[0]
    os.rename(extract_file, local_jsonfile)
    print(f"{host}: 下载完成{local_zipfile}, 解压到{local_jsonfile}")
    return local_jsonfile


def export_data_host(hostnames, dirpath="/opt/lavector/"):
    """
    导出所有标注完的数据
    :return:
    """
    print(time.strftime("%Y/%m/%d %H:%M:%S", time.localtime()))
    for hostname in hostnames:
        export_data(hostname=hostname, dirpath=dirpath)


def import_dev_data(hostname):
    """
    导入模型的开发数据集
    :return:
    """
    # testfile = "/Users/admin/git/TextBrewer/huazhuang/data_root_dir/newcos/dev.json"
    testfile = "/Users/admin/git/TextBrewer/huazhuang/utils/wrong.json"
    with open(testfile, 'r') as f:
        # 格式是[(text, keyword, labels)]
        loaddata = json.load(f)
    data = []
    for d in loaddata:
        # one_data = {'channel': 'jd','keyword': d[1],'text': d[0], 'wordtype': '未知'}
        one_data = {'data': {'channel': 'jd', 'keyword': d[1], 'text': d[0], 'wordtype': '未知',
                             "meta_info": {"location": "North Pole"}},
                    "completions": [
                        {
                            "result": [  # 标注结果, 这里对应的是一个人标注的结果，里面可能进行了多个标注
                                {
                                    "from_name": "label",
                                    "to_name": "text",
                                    "type": "labels",
                                    "value": {
                                        "end": d[3],
                                        "labels": [
                                            d[4]
                                        ],
                                        "start": d[2],
                                        "text": d[1]
                                    }
                                }
                            ]
                        }
                    ],
                    }
        data.append(one_data)
    r = requests.post(hostname + "project/import", data=json.dumps(data), headers=headers)
    pp.pprint(r.json())


def check_data(hostname):
    """
    查看下已导入的数据
    :return:
    """
    datas = get_tasks(hostname=hostname)
    not_repeat_id = []
    not_repeat_data = []
    for task in datas['tasks']:
        content = task['data']['text']
        if content not in not_repeat_data:
            not_repeat_data.append(content)
            not_repeat_id.append(task['id'])
        else:
            repeat_idx = not_repeat_data.index(content)
            repeat_id = not_repeat_id[repeat_idx]
            print(f"发现重复数据:{task['id']}和{repeat_id}")
    # print(f"共有重复数据{len(datas) - len(not_repeat_data)}条")
    return not_repeat_data, not_repeat_id

def import_pure_data(host, dirpath='/opt/lavector/absa', wordtype='包装', limit=None):
    """
    导入纯数据，不包含标签，从json文件中导入
    :param host: 主机列表
    :param dirpath: 文件夹或路径
    :param wordtype:      功效  成分  香味 包装 肤感
    :param limit:  限制导入数据条数，
    :return:
    """
    from convert_label_studio_data import get_all
    data = get_all(split=False, dirpath=dirpath)
    #过滤出是wordtype的数据
    filter_data = [d for d in data if d[6] == wordtype]
    valid_data = []
    unique_data = []
    for d in filter_data:
        if d[0]+ d[1] in unique_data:
            continue
        else:
            unique_data.append(d[0]+ d[1])
        one_data = {'channel': d[5], 'keyword': d[1], 'text': d[0], 'wordtype': d[-1]}
        valid_data.append(one_data)
    if limit and len(valid_data) > limit:
        valid_data = valid_data[:limit]
    every_host_number = int(len(valid_data) / len(host))
    print(f"每个主机导入数据{every_host_number}")
    vdatas = [valid_data[i:i + every_host_number] for i in range(0, len(valid_data), every_host_number)]
    for h, vdata in zip(host, vdatas):
        r = requests.post(h + "project/import", data=json.dumps(vdata), headers=headers)
        pp.pprint(r.json())
        print(f"共导入主机host{h}中数据{len(vdata)}条")


def import_excel_data(hostname):
    """
    导入模型人工标注后的excel模型, excel包含字段
    	Text	                Keyword	        Label	    Predict	    Location	    Probability	    Check
        控油效果：不错产品香味：很香泡沫数量：中适合发质：中性。。。	控油	中性	积极	(0, 2)	0.720	积极
    :return:
    """
    import pandas as pd
    testfile = "/Users/admin/Desktop/full_wrong.xlsx"
    df = pd.read_excel(testfile)
    data = []
    # 合并相同的Text，里面有n个关键字
    previous_keyword = None
    previsou_text = None
    result = []
    last_data = None
    for idx, d in df.iterrows():
        # one_data = {'channel': 'jd','keyword': d[1],'text': d[0], 'wordtype': '未知'}
        if d["Check"] == "无法判断":
            d["Check"] = "中性"
        # start和end必须是数字，否则无法显示
        one_result = {"from_name": "label",
                  "to_name": "text",
                  "type": "labels",
                  "value": {
                      "end": int(d["location"].lstrip('(').rstrip(')').split(',')[1]),
                      "labels": [
                          d["Check"]
                      ],
                      "start": int(d["location"].lstrip('(').rstrip(')').split(',')[0]),
                      "text": d["keyword"]
                  }
                  }
        # 初始化
        if previous_keyword is None and previsou_text is None:
            previous_keyword = d["keyword"]
            previsou_text = d["text"]
            last_data = {'data': {'channel': 'jd', 'keyword': d["keyword"], 'text': d["text"], 'wordtype': d["wordtype"],
                                 "meta_info": {"location": "North Pole"}},
                        "completions": [
                            {
                                "result": result  # 标注结果, 这里对应的是一个人标注的结果，里面可能进行了多个标注
                            }
                        ],
                        }
            result.append(one_result)
            continue
        elif previous_keyword == d["keyword"] and previsou_text == d["text"]:
            # 需要判断下一个句子的关键和text，如果相邻的2个句子和关键字是相同的，说明是一个句子中有几个关键词, 需要修改completions中的result
            result.append(one_result)
            #更新一下上一句的keyword
            previous_keyword = d["keyword"]
            previsou_text = d["text"]
            #更新last_data
            last_data = {'data': {'channel': 'jd', 'keyword': d["keyword"], 'text': d["text"], 'wordtype': d["wordtype"],
                                 "meta_info": {"location": "North Pole"}},
                        "completions": [
                            {
                                "result": result  # 标注结果, 这里对应的是一个人标注的结果，里面可能进行了多个标注
                            }
                        ],
                        }
            # 把下一句也添加进来, 如果不是最后一个元素，才继续循环下一句
            if idx == (len(df)-1):
                data.append(last_data)
        else:
            #把上一个词添加进去
            data.append(last_data)
            previous_keyword = d["keyword"]
            previsou_text = d["text"]
            #重置result,添加完成后
            result = []
            result.append(one_result)
            last_data = {'data': {'channel': 'jd', 'keyword': d["keyword"], 'text': d["text"], 'wordtype': d["wordtype"],
                                 "meta_info": {"location": "North Pole"}},
                        "completions": [
                            {
                                "result": result  # 标注结果, 这里对应的是一个人标注的结果，里面可能进行了多个标注
                            }
                        ],
                        }
            if idx == (len(df)-1):
                #最后一个元素也加进去
                data.append(last_data)
    r = requests.post(hostname + "project/import", data=json.dumps(data), headers=headers)
    pp.pprint(r.json())


def save_json_toexcel(jsonfile):
    """
    把导出的label-studio的导出的某个json文件导出为excel格式
    :param jsonfile: /opt/lavector/package/192.168.50.139_8081.json
    :return:
    """
    from convert_label_studio_data import get_all, save_excel
    data, cancel_data = get_all(keep_cancel=True, split=False, dirpath=jsonfile)
    save_excel(data=data,output_file='export.xlsx')

def read_all_text(ppt):
    prs = Presentation(ppt)
    # text_runs将被填入一个字符串列表，每个文本运行的演示文稿都有一个字符串。
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = ""
                for run in paragraph.runs:
                    paragraph_text = paragraph_text + run.text
                if paragraph_text:
                    text_runs.append(paragraph_text)

    return text_runs

def import_ppt_data(ppt,hostname):
    """
    从ppt中导入数据
    :param ppt:
    :type ppt:
    :return:
    :rtype:
    """
    data = []
    texts = read_all_text(ppt)
    not_repeat_data, not_repeat_id = check_data(hostname)
    # 过滤出中文来
    for text in texts:
        res = re.findall('[\u4e00-\u9fa5]+', text)
        if bool(res) and text not in not_repeat_id:
            data.append({'text': text})
    r = requests.post(hostname + "project/import", data=json.dumps(data), headers=headers)
    pp.pprint(r.json())
    print(f"共导入主机host{hostname}中数据{len(data)}条")


if __name__ == '__main__':
    # check_data()
    # setup_config(hostname=host)
    # get_project()
    # import_data()
    # get_tasks()
    # get_tasks(taskid=0)
    # get_completions()
    # delete_completions()
    # health()
    # list_models()
    # train_model()
    # predict_model()
    # hostnames = ["http://192.168.50.139:8081/api/"]
    # hostnames = ["http://192.168.50.139:8086/api/","http://192.168.50.139:8088/api/"]
    # hostnames = ["http://192.168.50.139:8081/api/","http://192.168.50.139:8082/api/", "http://192.168.50.139:8083/api/",
    #              "http://192.168.50.139:8084/api/","http://192.168.50.139:8085/api/","http://192.168.50.139:8086/api/",
    #              "http://192.168.50.139:8087/api/"]
    hostnames = ["http://127.0.0.1:8080/api/"]
    # setup_config(hostname="http://192.168.50.119:8090/api/")
    # import_absa_data_host(channel=['jd','tmall'],number=50, hostname=hostnames)
    # hostnames = ["http://192.168.50.119:8080/api/", "http://192.168.50.119:8081/api/"]
    # hostnames = ["http://192.168.50.119:8080/api/", "http://192.168.50.119:8081/api/","http://192.168.50.119:8082/api/",
    #              "http://192.168.50.119:8083/api/", "http://192.168.50.119:8084/api/","http://192.168.50.119:8085/api/",
    #              "http://192.168.50.119:8086/api/", "http://192.168.50.119:8087/api/","http://192.168.50.119:8088/api/",
    #              "http://192.168.50.119:8089/api/"]
    # import_absa_data_host_first(channel=['jd','tmall'],number=4000, hostname=hostnames)
    # get_tasks_host(hostnames=hostnames)
    import_ppt_data(ppt='/Users/admin/Downloads/翻译/【EL-国潮】中文版-0428.pptx',hostname=hostnames[0])
    # get_completions_host(hostnames=hostnames)
    # export_data(hostname="http://192.168.50.119:8090/api/")
    # export_data_host(hostnames=hostnames, dirpath="/opt/lavector/package/")
    # setup_config_host(hostnames=hostnames)
    # delete_tasks_host(hostnames=hostnames)
    # import_data(hostname=hostnames[0])
    # get_tasks(hostname='http://127.0.0.1:8080/api/')
    # ptimes1 = ["<:2020-10-01","<:2020-10-08", "<:2020-10-15","<:2020-10-30","<:2020-11-08","<:2020-11-15","<:2020-11-30","<:2020-12-08","<:2020-12-15", "<:2020-12-30", "<:2021-01-08","<:2021-1-15", "<:2021-1-30"]
    # ptimes2 = ["<:2020-09-01","<:2020-09-08", "<:2020-09-15","<:2020-09-20","<:2020-09-25","<:2020-11-11","<:2020-12-11","<:2020-12-25", "<:2021-01-08","<:2021-1-20", "<:2021-1-25"]
    # ptimes3 = ["<:2020-08-01","<:2020-08-15"]
    # ptimes = ptimes1+ ptimes2 +ptimes3
    # ptimes = ["<:2021-1-3", "<:2021-1-6"]
    # for ptime in ptimes:
    #     print(ptime)
    #     import_com_data_host_first(channel=None,keyword_threhold=0,ptime_keyword=ptime, leibie_num=[5000, -1, -1, -1, -1], require_tags=['component'], number=1000, unique_type=2, hostname=hostnames, not_cache=True, table="da_wide_table_new")
    # import_com_data_host_first(channel=None,keyword_threhold=20,ptime_keyword=">:2020-1-20", leibie_num=[5000, -1, -1, -1, -1], require_tags=['component'], number=5000, unique_type=2, hostname=hostnames, not_cache=True, table="da_wide_table_new")
    # import_dev_data(hostname=hostnames[0])
    # import_excel_data(hostname=hostnames[0])
    # import_pure_data(host=hostnames, wordtype='包装')
    # save_json_toexcel(jsonfile='/opt/lavector/package/192.168.50.139_8081.json')
